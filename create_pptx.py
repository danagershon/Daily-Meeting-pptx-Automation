from pptx import Presentation
from itertools import zip_longest
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.util import Pt

prs = Presentation('test.pptx')

finished_events = {"כותרת אירוע 1": "תיאור אירוע 1",
                   "כותרת אירוע 2": "תיאור אירוע 2"}


def print_shapes_names(slide=None):
    """
    print the shapes names in slide, or print the shapes of all slides if no slide is supplied.
    Every shape in a slide has a name that you can edit (in ppt or in the script).
    I have set the names of the shapes manually in ppt
    :param slide: the slide whose shapes should be printed.
    """
    if not slide:
        slides_to_iterate = prs.slides
    else:
        slides_to_iterate = [slide]

    for slide in slides_to_iterate:
        for shape in slide.shapes:
            print(f"name: {shape.name}")


def print_shapes_text(slide=None):
    """
    print the text of every shape in the slide, or print the text of each shape in all
    of the slides if no slide is supplied.
    Every shape in a slide can have a text_frame attribute, which has a text attribute that
    contains the text string
    :param slide: the slide whose shapes text should be printed.
    """
    if not slide:
        slides_to_iterate = prs.slides
    else:
        slides_to_iterate = [slide]

    for slide in slides_to_iterate:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            print(text_frame.text)


def fill_finished_events():
    """
    fills out the event header and description for each event in finished_events.
    It also takes care of formatting (font, color, size, ect.) with helper functions.
    In the finished event slide (slide number 1), each event has a shape group with a
    header shape (shape 0 in the group) and a description shape (shape 1 in the group)
    """
    slide = prs.slides[1]
    for event_group, (event_header, event_desc) in zip_longest(slide.shapes, finished_events.items(), fillvalue=("", "")):
        event_header_shape = event_group.shapes[0]
        format_event_header(event_header_shape.text_frame, event_header)

        event_desc_shape = event_group.shapes[1]
        format_event_desc(event_desc_shape.text_frame, event_desc)


def format_event_header(event_header_tf, text=""):
    event_header_tf.clear()  # removes all paragraphs in text frame but keeps a default empty paragraph
    header_p = event_header_tf.paragraphs[0]
    header_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    """ 
    in order to format the text, you have to create a "run" object in the paragraph, 
    which has the attributes for formatting such as size, font, ect. 
    """
    header_run = header_p.add_run()
    header_run.text = text

    header_font = header_run.font
    header_font.color.rgb = RGBColor(0, 0, 0)  # black
    header_font.bold = True
    header_font.size = Pt(24)
    header_font.name = 'Calibri'  # change to Heebo later :)


def format_event_desc(event_desc_tf, text=""):
    event_desc_tf.clear()
    desc_p = event_desc_tf.paragraphs[0]
    desc_p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

    desc_run = desc_p.add_run()
    desc_run.text = text

    desc_font = desc_run.font
    desc_font.color.rgb = RGBColor(0, 0, 0)
    desc_font.size = Pt(18)
    desc_font.name = 'Calibri'


fill_finished_events()
print_shapes_names(prs.slides[1])
prs.save('new.pptx')
