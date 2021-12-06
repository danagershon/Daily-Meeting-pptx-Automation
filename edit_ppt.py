from pptx import Presentation

prs = Presentation('test.pptx')

for slide in prs.slides:
    for i, shape in enumerate(slide.shapes):
        print("id: %s, type: %s" % (shape.name, shape.shape_type))
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        text_frame.text = str(i)

prs.save('new.pptx')