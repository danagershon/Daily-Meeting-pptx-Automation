from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Cm
from pptx.util import Inches

prs = Presentation('test.pptx')
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes
lefts = [Cm(0.32), Cm(8.72), Cm(17.11), Cm(25.51)]
top = Cm(3)
width = Cm(8)
height = Cm(10)

for i in range(4):
    shape = shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, lefts[i], top, width, height
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(197, 224, 180)
    line = shape.line
    line.color.rgb = RGBColor(0, 0, 0)

prs.save('test.pptx')
